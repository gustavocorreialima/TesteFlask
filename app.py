from flask import Flask, request, send_file, jsonify
import json
import os
from pptx import Presentation
from datetime import datetime

app = Flask(__name__)

# Configurações
TEMPLATE_PATH = "template.pptx"
UPLOAD_FOLDER = "temp"

# Cria pasta temporária se não existir
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def preencher_pptx(template_path, dados, output_path):
    """Função adaptada do seu script"""
    prs = Presentation(template_path)

    # Se vier uma lista com vários objetos, pegue o primeiro
    if isinstance(dados, list) and len(dados) > 0:
        dados = dados[0]["response"]["output"]

    # Mapeamento manual para corrigir diferença check_pp → check_PP
    if "check_pp" in dados:
        dados["check_PP"] = dados["check_pp"]

    # Itera pelos slides do PPTX
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texto_original = shape.text
                novo_texto = texto_original

                # Substitui cada {{campo}} presente no template
                for chave, valor in dados.items():
                    placeholder = f"{{{{{chave}}}}}"
                    novo_texto = novo_texto.replace(placeholder, str(valor))

                shape.text = novo_texto

    # Salva arquivo final
    prs.save(output_path)
    return output_path


@app.route('/gerar-pptx', methods=['POST'])
def gerar_pptx():
    try:
        # Recebe o JSON da requisição
        dados_json = request.get_json()
        
        if not dados_json:
            return jsonify({'erro': 'Nenhum JSON foi enviado'}), 400
        
        # Gera nome único para o arquivo de saída
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_filename = f"business_case_{timestamp}.pptx"
        output_path = os.path.join(UPLOAD_FOLDER, output_filename)
        
        # Processa o PPTX
        preencher_pptx(TEMPLATE_PATH, dados_json, output_path)
        
        # Retorna o arquivo gerado
        return send_file(
            output_path,
            as_attachment=True,
            download_name=output_filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        
    except Exception as e:
        return jsonify({'erro': str(e)}), 500


@app.route('/health', methods=['GET'])
def health():
    """Rota para verificar se a API está funcionando"""
    return jsonify({'status': 'OK', 'mensagem': 'API funcionando'}), 200


if __name__ == '__main__':
    app.run(debug=False, host='0.0.0.0', port=5000)